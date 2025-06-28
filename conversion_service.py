# conversion_service.py
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import JSONResponse
import subprocess
import io
import base64
import fitz # PyMuPDF
import os
import shutil
import tempfile
import json # For logging output

app = FastAPI()

# --- Helper functions for document conversion ---

def _convert_pptx_to_images_and_text(pptx_bytes: bytes) -> list[dict]:
    """
    Converts a PPTX file to a list of Base64 encoded PNG images, one per slide,
    and extracts text from each slide. Requires LibreOffice and unoconv.
    """
    results = []
    
    # Create a temporary directory for safe file handling
    with tempfile.TemporaryDirectory() as temp_dir:
        pptx_path = os.path.join(temp_dir, "input.pptx")
        output_dir = os.path.join(temp_dir, "output_images")
        os.makedirs(output_dir)

        # Write the uploaded PPTX bytes to a temporary file
        with open(pptx_path, "wb") as f:
            f.write(pptx_bytes)

        try:
            # Command to convert slides to individual PNGs
            command_render = [
                "unoconv", "-f", "png", 
                "--output", os.path.join(output_dir, "slide-.png"), # unoconv adds page number
                pptx_path
            ]
            # MODIFIED: Capture output and check return code manually for better error reporting
            result = subprocess.run(command_render, capture_output=True, text=True, check=False)
            if result.returncode != 0:
                # Raise a CalledProcessError to be caught below, including stdout/stderr
                raise subprocess.CalledProcessError(result.returncode, command_render, output=result.stdout, stderr=result.stderr)
            
            # Use python-pptx to extract text, as unoconv doesn't give text directly
            from pptx import Presentation
            prs = Presentation(io.BytesIO(pptx_bytes))

            for i, slide in enumerate(prs.slides):
                slide_text_content = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_text_content.append(shape.text)
                text = " ".join(slide_text_content)[:2000] # Limit text length

                # Find the corresponding image file
                # unoconv creates files like slide-1.png, slide-2.png etc.
                image_path = os.path.join(output_dir, f"slide-{i+1}.png") 
                if not os.path.exists(image_path):
                    # Fallback for some unoconv versions or edge cases, try alternative naming
                    image_path = os.path.join(output_dir, f"slide-{i+1:02d}.png") # e.g., slide-01.png
                if not os.path.exists(image_path):
                    # If still not found, it might be an issue with unoconv or indexing.
                    # Provide a warning and continue without image for this slide.
                    print(f"Warning: Image file not found for slide {i+1} at {image_path}")
                    image_data = "" # No image data
                else:
                    with open(image_path, "rb") as img_file:
                        image_data = base64.b64encode(img_file.read()).decode('utf-8')

                results.append({
                    "slide_index": i,
                    "text": text,
                    "image_data": image_data
                })
        except subprocess.CalledProcessError as e:
            # MODIFIED: Include both stdout and stderr in the error detail
            detail_message = f"PPTX conversion failed. unoconv exited with code {e.returncode}.\n"
            if e.stdout:
                detail_message += f"stdout:\n{e.stdout}\n"
            if e.stderr:
                detail_message += f"stderr:\n{e.stderr}\n"
            if not e.stdout and not e.stderr:
                detail_message += "No output captured from unoconv (check if LibreOffice is fully configured for headless mode)."
            raise HTTPException(status_code=500, detail=detail_message)
        except FileNotFoundError:
            raise HTTPException(status_code=500, detail="unoconv or LibreOffice not found. Please ensure LibreOffice and unoconv are installed and in PATH.")
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"An unexpected error occurred during PPTX processing: {e}")
            
    return results

def _convert_pdf_to_images_and_text(pdf_bytes: bytes) -> list[dict]:
    """
    Converts a PDF file to a list of Base64 encoded PNG images, one per page,
    and extracts text from each page. Uses PyMuPDF (fitz).
    """
    results = []
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        for i, page in enumerate(doc):
            # Render page to PNG image
            pix = page.get_pixmap()
            image_bytes = pix.tobytes("png")
            image_data = base64.b64encode(image_bytes).decode('utf-8')

            # Extract text from page
            text = page.get_text()[:2000] # Limit text length

            results.append({
                "slide_index": i, # Using slide_index for consistency with PPTX
                "text": text,
                "image_data": image_data
            })
        doc.close()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF conversion failed: {e}")
    return results

@app.post("/convert_document")
async def convert_document_endpoint(file: UploadFile = File(...)):
    """
    Endpoint to convert an uploaded PPTX or PDF document into a list of
    slide/page data (text and Base64 image).
    """
    file_bytes = await file.read()
    file_type = file.content_type

    if file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        slides_data = _convert_pptx_to_images_and_text(file_bytes)
    elif file_type == 'application/pdf':
        slides_data = _convert_pdf_to_images_and_text(file_bytes)
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {file_type}. Only PPTX and PDF are supported.")
    
    return JSONResponse(content={"slides": slides_data})

import streamlit as st
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.dml.color import RGBColor
import io
import copy
import uuid
import openai
import json
import requests
import os
import subprocess
import tempfile
import shutil
import mimetypes
import httpx # <-- ADD THIS IMPORT

# --- Configuration for the Conversion Service ---
CONVERSION_SERVICE_URL = os.getenv("CONVERSION_SERVICE_URL", "http://localhost:8000/convert_document")

# --- Helper Function for Copying Background (PPTX-specific) ---
def copy_slide_background(src_slide, dest_slide):
    # This function remains unchanged
    src_slide_elm = src_slide.element
    dest_slide_elm = dest_slide.element
    src_bg_pr = src_slide_elm.find('.//p:bgPr', namespaces=src_slide_elm.nsmap)
    if src_bg_pr is None:
        return
    src_blip_fill = src_bg_pr.find('.//a:blipFill', namespaces=src_slide_elm.nsmap)
    if src_blip_fill is not None:
        src_blip = src_blip_fill.find('.//a:blip', namespaces=src_blip_fill.nsmap)
        if src_blip is not None and '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed' in src_blip.attrib:
            rId = src_blip.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed']
            try:
                src_image_part = src_slide.part.related_part(rId)
                image_bytes = src_image_part.blob
                new_image_part = dest_slide.part.get_or_add_image_part(image_bytes, src_image_part.content_type)
                new_rId = dest_slide.part.relate_to(new_image_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image')
                new_bg_pr = copy.deepcopy(src_bg_pr)
                new_blip = new_bg_pr.find('.//a:blip', namespaces=new_bg_pr.nsmap)
                if new_blip is not None:
                    new_blip.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'] = new_rId
                current_bg = dest_slide_elm.find('.//p:bg', namespaces=dest_slide_elm.nsmap)
                if current_bg is not None:
                    current_bg.getparent().remove(current_bg)
                dest_slide_elm.append(new_bg_pr)
            except Exception as e:
                print(f"Warning: Could not copy background image. Error: {e}")
                copy_solid_or_gradient_background(src_slide, dest_slide)
    else:
        copy_solid_or_gradient_background(src_slide, dest_slide)

def copy_solid_or_gradient_background(src_slide, dest_slide):
    # This function remains unchanged
    src_slide_elm = src_slide.element
    dest_slide_elm = dest_slide.element
    src_bg_pr = src_slide_elm.find('.//p:bgPr', namespaces=src_slide_elm.nsmap)
    if src_bg_pr is not None:
        new_bg_pr = copy.deepcopy(src_bg_pr)
        current_bg = dest_slide_elm.find('.//p:bg', namespaces=dest_slide_elm.nsmap)
        if current_bg is not None:
            current_bg.getparent().remove(current_bg)
        dest_slide_elm.append(new_bg_pr)

# --- Core PowerPoint Functions (for PPTX output generation) ---
def deep_copy_slide_content(dest_slide, src_slide):
    # This function remains unchanged
    for shape in list(dest_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)
    for shape in src_slide.shapes:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image_bytes = shape.image.blob
                dest_slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width, height)
            except Exception as e:
                print(f"Warning: Could not copy picture from source slide. Error: {e}")
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    new_el = copy.deepcopy(shape.element)
                    dest_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        elif shape.has_text_frame:
            new_shape = dest_slide.shapes.add_textbox(left, top, width, height)
            new_text_frame = new_shape.text_frame
            new_text_frame.clear()
            for paragraph in shape.text_frame.paragraphs:
                new_paragraph = new_text_frame.add_paragraph()
                new_paragraph.alignment = paragraph.alignment
                if hasattr(paragraph, 'level'):
                    new_paragraph.level = paragraph.level
                for run in paragraph.runs:
                    new_run = new_paragraph.add_run()
                    new_run.text = run.text
                    new_run.font.bold = run.font.bold
                    new_run.font.italic = run.font.italic
                    new_run.font.underline = run.font.underline
                    if run.font.size:
                        new_run.font.size = run.font.size
                    if run.font.fill.type == MSO_FILL_TYPE.SOLID:
                        new_run.font.fill.solid()
                        try:
                            if isinstance(run.font.fill.fore_color.rgb, RGBColor):
                                new_run.font.fill.fore_color.rgb = run.font.fill.fore_color.rgb
                            else: 
                                rgb_tuple = run.font.fill.fore_color.rgb
                                new_run.font.fill.fore_color.rgb = RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])
                        except Exception as color_e:
                            print(f"Warning: Could not copy font color. Error: {color_e}")
                            pass
            new_text_frame.word_wrap = shape.text_frame.word_wrap
            new_text_frame.margin_left = shape.text_frame.margin_left
            new_text_frame.margin_right = shape.text_frame.margin_right
            new_text_frame.margin_top = shape.text_frame.margin_top
            new_text_frame.margin_bottom = shape.text_frame.margin_bottom
        else:
            new_el = copy.deepcopy(shape.element)
            dest_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    copy_slide_background(src_slide, dest_slide)

def get_all_slide_data(file_bytes: bytes, file_type: str) -> list[dict]:
    files = {'file': (f"document.{file_type.split('/')[-1]}", file_bytes, file_type)}
    try:
        response = requests.post(CONVERSION_SERVICE_URL, files=files, timeout=300)
        response.raise_for_status()
        return response.json()['slides']
    except requests.exceptions.RequestException as e:
        st.error(f"Error connecting to conversion service or during conversion: {e}. Please ensure the conversion service is running at {CONVERSION_SERVICE_URL} and has required system dependencies (LibreOffice, PyMuPDF).")
        st.stop()
    except KeyError:
        st.error("Conversion service returned an unexpected response format.")
        st.stop()

def find_slide_by_ai(api_key, file_bytes: bytes, file_type: str, slide_type_prompt: str, deck_name: str):
    if not slide_type_prompt: return {"slide": None, "index": -1, "justification": "No keyword provided."}
    if not api_key:
        return {"slide": None, "index": -1, "justification": "OpenAI API Key is missing."}

    # --- START OF THE FIX ---
    # Manually create an HTTP client with proxies explicitly disabled
    http_client = httpx.Client(proxies={})
    client = openai.OpenAI(api_key=api_key, http_client=http_client)
    # --- END OF THE FIX ---
    
    slides_data = get_all_slide_data(file_bytes, file_type)
    system_prompt = f"""
    You are an expert presentation analyst. Your task is to find the best slide/page in a document that matches a user's description.
    The user is looking for a slide/page representing: '{slide_type_prompt}'.
    Analyze both the provided **text content** and the **visual structure (from the image)** for each slide/page to infer its purpose.
    For 'Timeline' slides/pages: Look for strong textual indicators of sequential progression and visual patterns that imply a timeline.
    For 'Objectives' slides/pages: These will typically contain goal-oriented language.
    You must prioritize actual content slides/pages over simple divider or table of contents pages.
    Return a JSON object with 'best_match_index' (integer, or -1) and 'justification' (brief, one-sentence).
    """
    user_parts = [
        {"type": "text", "text": f"Find the best slide/page for '{slide_type_prompt}' in the '{deck_name}' with the following pages/slides:"}
    ]
    for slide_info in slides_data:
        user_parts.append({"type": "text", "text": f"\n--- Page/Slide {slide_info['slide_index'] + 1} (Text): {slide_info['text']}"})
        user_parts.append({
            "type": "image_url",
            "image_url": { "url": f"data:image/png;base64,{slide_info['image_data']}" }
        })
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_parts}
    ]
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=messages,
            response_format={"type": "json_object"}
        )
        result = json.loads(response.choices[0].message.content)
        best_index = result.get("best_match_index", -1)
        justification = result.get("justification", "No justification provided.")
        selected_slide_data = slides_data[best_index] if best_index != -1 and best_index < len(slides_data) else None
        return {"slide": selected_slide_data, "index": best_index, "justification": justification}
    except openai.APIError as e:
        return {"slide": None, "index": -1, "justification": f"OpenAI API Error: {e}"}
    except json.JSONDecodeError as e:
        return {"slide": None, "index": -1, "justification": f"AI response was not valid JSON: {e}"}
    except Exception as e:
        return {"slide": None, "index": -1, "justification": f"An unexpected error occurred during AI analysis: {e}"}

def analyze_and_map_content(api_key, gtm_slide_content_data, template_slides_data, user_keyword):
    if not api_key:
        return {"best_template_index": -1, "justification": "OpenAI API Key is missing.", "processed_content": gtm_slide_content_data}

    # --- START OF THE FIX ---
    # Manually create an HTTP client with proxies explicitly disabled
    http_client = httpx.Client(proxies={})
    client = openai.OpenAI(api_key=api_key, http_client=http_client)
    # --- END OF THE FIX ---

    system_prompt = f"""
    You are an expert presentation content mapper. Your primary task is to help a user integrate content from a Global (GTM) slide/page into the most appropriate regional template.
    Given the `gtm_slide_content` (with its text and image) and a list of `template_slides_data` (each with an index and text content, and image data), you must perform two critical tasks:
    1. Select the BEST Template by visually and semantically evaluating which template slide's structure would best accommodate the GTM content.
    2. Process GTM Content for Regionalization: Analyze the `gtm_slide_content` and replace any regional-specific parts with a generic placeholder like `[REGIONAL DATA HERE]`.
    You MUST return a JSON object with 'best_template_index' (integer), 'justification' (string), and 'processed_content' (object with 'title' and 'body').
    """
    user_parts = [
        {"type": "text", "text": f"User's original keyword for this content: '{user_keyword}'"},
        {"type": "text", "text": "GTM Slide/Page Content to Process (Text):"},
        {"type": "text", "text": json.dumps(gtm_slide_content_data.get('text', {}), indent=2)},
        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{gtm_slide_content_data['image_data']}"}} 
    ]
    user_parts.append({"type": "text", "text": "\nAvailable Template Slides/Pages Summary and Visuals:"})
    for slide_info in template_slides_data:
        user_parts.append({"type": "text", "text": f"\n--- Template Slide/Page {slide_info['slide_index'] + 1} (Text): {slide_info['text']}"})
        user_parts.append({
            "type": "image_url",
            "image_url": { "url": f"data:image/png;base64,{slide_info['image_data']}" }
        })
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_parts}
    ]
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=messages,
            response_format={"type": "json_object"}
        )
        result = json.loads(response.choices[0].message.content)
        if "best_template_index" not in result or "justification" not in result or "processed_content" not in result:
            raise ValueError("AI response missing required keys.")
        best_index = result["best_template_index"]
        justification = result["justification"]
        processed_content = result["processed_content"]
        if "title" not in processed_content: processed_content["title"] = gtm_slide_content_data.get("title", "")
        if "body" not in processed_content: processed_content["body"] = gtm_slide_content_data.get("body", "")
        return {"best_template_index": best_index, "justification": justification, "processed_content": processed_content}
    except openai.APIError as e:
        print(f"OpenAI API Error in analyze_and_map_content: {e}")
        return {"best_template_index": -1, "justification": f"OpenAI API Error: {e}", "processed_content": gtm_slide_content_data}
    except json.JSONDecodeError as e:
        print(f"JSON Decode Error in analyze_and_map_content: {e}")
        return {"best_template_index": -1, "justification": f"AI response was not valid JSON: {e}", "processed_content": gtm_slide_content_data}
    except Exception as e:
        print(f"An unexpected error occurred in analyze_and_map_content: {e}")
        return {"best_template_index": -1, "justification": f"An error occurred during content mapping: {e}", "processed_content": gtm_slide_content_data}

def get_slide_content(slide):
    # This function remains unchanged
    if not slide: return {"title": "", "body": ""}
    text_shapes = sorted([s for s in slide.shapes if s.has_text_frame and s.text.strip()], key=lambda s: s.top)
    title = ""
    body = ""
    if text_shapes:
        title = text_shapes[0].text.strip()
        body = "\n".join(s.text.strip() for s in text_shapes[1:])
    return {"title": title, "body": body}

def populate_slide(slide, content):
    # This function remains unchanged
    title_populated, body_populated = False, False
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        is_title_placeholder = (hasattr(shape, 'is_placeholder') and shape.is_placeholder and shape.placeholder_format.type in (1, 2, 8))
        is_top_text_box = (shape.top < Pt(150))
        if not title_populated and (is_title_placeholder or is_top_text_box):
            tf = shape.text_frame
            tf.clear()
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = content.get("title", "")
            title_populated = True
        is_body_placeholder = (hasattr(shape, 'is_placeholder') and shape.is_placeholder and shape.placeholder_format.type in (3, 4, 8, 14))
        is_lorem_ipsum = "lorem ipsum" in shape.text.lower()
        is_empty_text_box = not shape.text.strip() and shape.height > Pt(100)
        if not body_populated and (is_body_placeholder or is_lorem_ipsum or is_empty_text_box):
            tf = shape.text_frame
            tf.clear()
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = content.get("body", "")
            body_populated = True
        if title_populated and body_populated:
            break

# --- Streamlit App ---
st.set_page_config(page_title="Dynamic AI Presentation Assembler", layout="wide")
st.title("📊 Dynamic AI Presentation Assembler")

with st.sidebar:
    st.header("1. API Key")
    api_key = st.text_input("OpenAI API Key", type="password")
    st.markdown("---")
    st.header("2. Input Documents (Drag & Drop)")
    st.info("Upload your PPTX or PDF files directly.")
    st.subheader("Upload Template Documents")
    uploaded_template_files = st.file_uploader(
        "Drag and drop your Template PPTX/PDF files here",
        type=["pptx", "pdf"],
        accept_multiple_files=True,
        key="template_uploader"
    )
    st.subheader("Upload GTM Global Document")
    uploaded_gtm_file = st.file_uploader(
        "Drag and drop your GTM Global PPTX/PDF file here",
        type=["pptx", "pdf"],
        accept_multiple_files=False,
        key="gtm_uploader"
    )
    st.markdown("---")
    st.header("3. Define Presentation Structure")
    if 'structure' not in st.session_state: 
        st.session_state.structure = []
    if st.button("Add New Step", use_container_width=True):
        st.session_state.structure.append({"id": str(uuid.uuid4()), "keyword": "", "action": "Copy from GTM (as is)"})
    for i, step in enumerate(st.session_state.structure):
        with st.container(border=True):
            cols = st.columns([3, 3, 1])
            step["keyword"] = cols[0].text_input("Slide Type", step["keyword"], key=f"keyword_{step['id']}")
            step["action"] = cols[1].selectbox(
                "Action", 
                ["Copy from GTM (as is)", "Merge: Template Layout + GTM Content"], 
                index=["Copy from GTM (as is)", "Merge: Template Layout + GTM Content"].index(step["action"]), 
                key=f"action_{step['id']}"
            )
            if cols[2].button("🗑️", key=f"del_{step['id']}"):
                st.session_state.structure.pop(i)
                st.rerun()
    if st.button("Clear Structure", use_container_width=True): 
        st.session_state.structure = []
        st.rerun()

# --- Main App Logic ---
if uploaded_template_files and uploaded_gtm_file and api_key and st.session_state.structure:
    if st.button("🚀 Assemble Presentation", type="primary"):
        with st.spinner("Assembling your new presentation..."):
            try:
                st.write("Step 1/3: Loading and processing uploaded documents...")
                all_template_slides_for_ai = []
                base_pptx_template_found = False
                new_prs = None 
                for uploaded_file in uploaded_template_files:
                    file_bytes = uploaded_file.read()
                    file_type = uploaded_file.type
                    file_name = uploaded_file.name
                    if file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                        if not base_pptx_template_found:
                            new_prs = Presentation(io.BytesIO(file_bytes))
                            st.info(f"Using '{file_name}' as the primary base PPTX template.")
                            base_pptx_template_found = True
                        else:
                            current_prs_to_merge = Presentation(io.BytesIO(file_bytes))
                            st.info(f"Merging slides from '{file_name}' into the base template.")
                            for slide_to_merge in current_prs_to_merge.slides:
                                new_slide = new_prs.slides.add_slide(new_prs.slide_layouts[0]) 
                                deep_copy_slide_content(new_slide, slide_to_merge) 
                    all_template_slides_for_ai.extend(get_all_slide_data(file_bytes, file_type))
                if new_prs is None:
                    st.error("Error: At least one PPTX file must be uploaded as a 'Template Document' to serve as the base for the assembled presentation.")
                    st.stop() 
                gtm_file_to_process_bytes = uploaded_gtm_file.read()
                gtm_file_to_process_type = uploaded_gtm_file.type
                gtm_file_to_process_name = uploaded_gtm_file.name
                st.info(f"Using '{gtm_file_to_process_name}' as the GTM Global Document.")
                process_log = []
                st.write("Step 2/3: Building new presentation based on your structure...")
                num_template_slides = len(new_prs.slides) 
                num_structure_steps = len(st.session_state.structure)
                if num_structure_steps < num_template_slides:
                    for i in range(num_template_slides - 1, num_structure_steps - 1, -1):
                        rId = new_prs.slides._sldIdLst[i].rId
                        new_prs.part.drop_rel(rId)
                        del new_prs.slides._sldIdLst[i]
                    st.info(f"Removed {num_template_slides - num_structure_steps} unused slides from the merged template.")
                elif num_structure_steps > num_template_slides:
                     st.warning(f"Warning: Your defined structure has more steps ({num_structure_steps}) than the merged template has slides ({num_template_slides}). Extra steps will be ignored.")

                for i, step in enumerate(st.session_state.structure):
                    if i >= len(new_prs.slides): 
                        break

                    current_dest_slide_index = i
                    dest_slide = new_prs.slides[current_dest_slide_index] 
                    keyword = step["keyword"]
                    action = step["action"]
                    log_entry = {"step": i + 1, "keyword": keyword, "action": action, "log": []}
                    
                    if action == "Copy from GTM (as is)":
                        if gtm_file_to_process_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation': 
                            gtm_prs = Presentation(io.BytesIO(gtm_file_to_process_bytes))
                            result = find_slide_by_ai(api_key, gtm_file_to_process_bytes, gtm_file_to_process_type, keyword, "GTM Deck")
                            log_entry["log"].append(f"**GTM Content Choice Justification (PPTX Copy):** {result['justification']}")
                            if result["slide"]:
                                src_slide_object = gtm_prs.slides[result["index"]] 
                                deep_copy_slide_content(dest_slide, src_slide_object)
                                log_entry["log"].append(f"**Action:** Replaced Template slide {current_dest_slide_index + 1} with content from GTM PPTX slide {result['index'] + 1}.")
                            else:
                                log_entry["log"].append("**Action:** No suitable slide found in GTM PPTX deck. Template slide was left as is.")
                        else: 
                            log_entry["log"].append(f"**Warning:** 'Copy from GTM (as is)' is selected but GTM deck is a PDF. This action cannot directly copy PPTX shapes from a PDF. Proceeding with 'Merge' logic for content extraction based on text and assumed visuals.")
                            gtm_ai_selection_result = find_slide_by_ai(api_key, gtm_file_to_process_bytes, gtm_file_to_process_type, keyword, "GTM Deck (Content Source)")
                            log_entry["log"].append(f"**GTM Content Source Justification (PDF Fallback Merge):** {gtm_ai_selection_result['justification']}")
                            raw_gtm_content = {"title": "", "body": ""}
                            if gtm_ai_selection_result["slide"]:
                                full_text = gtm_ai_selection_result["slide"].get("text", "")
                                lines = full_text.split('\n')
                                raw_gtm_content["title"] = lines[0] if lines else ""
                                raw_gtm_content["body"] = "\n".join(lines[1:]) if len(lines) > 1 else ""
                            ai_mapping_result = analyze_and_map_content(api_key, raw_gtm_content, all_template_slides_for_ai, keyword)
                            log_entry["log"].append(f"**AI Template Mapping Justification (PDF Fallback Merge):** {ai_mapping_result['justification']}")
                            selected_template_index = ai_mapping_result["best_template_index"]
                            processed_content = ai_mapping_result["processed_content"]
                            if selected_template_index != -1 and selected_template_index < len(new_prs.slides):
                                populate_slide(dest_slide, processed_content)
                                log_entry["log"].append(f"**Action:** Merged processed content from GTM (PDF) page {gtm_ai_selection_result['index'] + 1} into Template slide {current_dest_slide_index + 1}, with regional placeholders. AI suggested template type from template index {selected_template_index + 1}.")
                            else:
                                log_entry["log"].append("**Action:** AI could not determine a suitable template layout or process content for PDF. Template slide was left as is.")

                    elif action == "Merge: Template Layout + GTM Content":
                        gtm_ai_selection_result = find_slide_by_ai(api_key, gtm_file_to_process_bytes, gtm_file_to_process_type, keyword, "GTM Deck (Content Source)")
                        log_entry["log"].append(f"**GTM Content Source Justification:** {gtm_ai_selection_result['justification']}")
                        raw_gtm_content = {"title": "", "body": ""}
                        if gtm_ai_selection_result["slide"]:
                            full_text = gtm_ai_selection_result["slide"].get("text", "")
                            lines = full_text.split('\n')
                            raw_gtm_content["title"] = lines[0] if lines else ""
                            raw_gtm_content["body"] = "\n".join(lines[1:]) if len(lines) > 1 else ""
                        ai_mapping_result = analyze_and_map_content(api_key, raw_gtm_content, all_template_slides_for_ai, keyword)
                        log_entry["log"].append(f"**AI Template Mapping Justification:** {ai_mapping_result['justification']}")
                        selected_template_index = ai_mapping_result["best_template_index"]
                        processed_content = ai_mapping_result["processed_content"]
                        if selected_template_index != -1 and selected_template_index < len(new_prs.slides):
                            populate_slide(dest_slide, processed_content)
                            log_entry["log"].append(f"**Action:** Merged processed content from GTM ({gtm_file_to_process_name}) page/slide {gtm_ai_selection_result['index'] + 1} into Template slide {current_dest_slide_index + 1}, with regional placeholders. AI suggested template type from template index {selected_template_index + 1}.")
                        else:
                            log_entry["log"].append("**Action:** AI could not determine a suitable template layout or process content. Template slide was left as is.")
                    process_log.append(log_entry)
 
                st.success("Successfully built the new presentation structure.")
                st.write("Step 3/3: Finalizing...")
                st.subheader("📋 Process Log")
                for entry in process_log:
                    with st.expander(f"Step {entry['step']}: '{entry['keyword']}' ({entry['action']})"):
                        for line in entry['log']: 
                            st.markdown(f"- {line}")
                output_buffer = io.BytesIO()
                new_prs.save(output_buffer)
                output_buffer.seek(0)
                st.success("✨ Your new regional presentation has been assembled!")
                st.download_button(
                    "Download Assembled PowerPoint", 
                    data=output_buffer, 
                    file_name="Dynamic_AI_Assembled_Deck.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception as e:
                st.error(f"A critical error occurred: {e}")
                st.exception(e)
else:
    st.info("Please provide an API Key, upload your Template/GTM documents, and define the structure in the sidebar to begin.")
